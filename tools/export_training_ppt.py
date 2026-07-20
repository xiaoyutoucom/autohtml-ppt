"""将培训 HTML 逐页截取浏览器可视区域，合成 PPTX。

依赖：
  pip install -r requirements.txt
  playwright install chromium

用法（本工程根目录 E:\\wence\\autohtml-ppt）：
  python -m tools.export_training_ppt
  python -m tools.export_training_ppt --html docs/harness_training.html --out docs/Harness_Engineering_培训.pptx
"""

from __future__ import annotations

import argparse
import asyncio
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]

FIT_SLIDE_JS = """
() => {
  const root = document.documentElement;
  root.classList.add('exporting-ppt');
  const footer = document.querySelector('.footer');
  if (footer) footer.classList.remove('is-peek');

  const s = document.querySelector('.slide.active');
  if (!s) return { ok: false };

  s.style.animation = 'none';
  s.style.opacity = '1';
  s.style.filter = 'none';
  s.style.transform = 'none';
  s.style.transformOrigin = 'center center';

  // 用子元素包围盒判断是否超出视口（封面 cover-foot 常被裁）
  const host = s.getBoundingClientRect();
  let maxBottom = 0;
  let maxRight = 0;
  Array.from(s.children).forEach((k) => {
    const r = k.getBoundingClientRect();
    maxBottom = Math.max(maxBottom, r.bottom - host.top);
    maxRight = Math.max(maxRight, r.right - host.left);
  });
  // 再扫一层深层可见块，避免 grid 子项未包住全部
  s.querySelectorAll('.cover-foot, .cover-hero, .cover-speaker, .ref-panel, table.fw-table, .takeaway-grid').forEach((k) => {
    const r = k.getBoundingClientRect();
    maxBottom = Math.max(maxBottom, r.bottom - host.top);
    maxRight = Math.max(maxRight, r.right - host.left);
  });

  const needH = Math.max(maxBottom, host.height);
  const needW = Math.max(maxRight, host.width);
  let scale = Math.min(1, host.height / needH, host.width / needW);
  if (scale < 0.999) {
    scale = Math.max(0.55, scale * 0.96);
    s.style.transform = `scale(${scale})`;
  }
  return { ok: true, scale, needH, hostH: host.height };
}
"""


def _parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Export training deck HTML slides to PPTX (viewport screenshots).")
    p.add_argument(
        "--html",
        type=Path,
        default=REPO / "docs" / "harness_training.html",
        help="Training HTML path",
    )
    p.add_argument(
        "--out",
        type=Path,
        default=REPO / "docs" / "Harness_Engineering_培训.pptx",
        help="Output PPTX path",
    )
    p.add_argument("--width", type=int, default=1920)
    p.add_argument("--height", type=int, default=1080)
    p.add_argument("--settle-ms", type=int, default=550, help="Wait after each page change")
    return p.parse_args()


async def _export(args: argparse.Namespace) -> Path:
    try:
        from playwright.async_api import async_playwright
    except ImportError as e:
        raise SystemExit("需要 playwright：pip install playwright && playwright install chromium") from e
    try:
        from pptx import Presentation
        from pptx.util import Emu, Inches
    except ImportError as e:
        raise SystemExit("需要 python-pptx：pip install python-pptx") from e

    html_path = args.html.resolve()
    if not html_path.is_file():
        raise SystemExit(f"HTML 不存在: {html_path}")

    out_path = args.out.resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    uri = html_path.as_uri()

    shots: list[Path] = []
    with tempfile.TemporaryDirectory(prefix="deck_ppt_") as tmp:
        tmp_dir = Path(tmp)
        async with async_playwright() as pw:
            browser = await pw.chromium.launch(headless=True)
            page = await browser.new_page(
                viewport={"width": args.width, "height": args.height},
                device_scale_factor=1,
            )
            await page.goto(uri, wait_until="domcontentloaded")
            try:
                await page.evaluate("() => document.fonts && document.fonts.ready")
            except Exception:
                pass
            await page.wait_for_timeout(900)

            await page.evaluate(
                """() => {
                  document.documentElement.classList.add('exporting-ppt');
                  const f = document.querySelector('.footer');
                  if (f) f.classList.remove('is-peek');
                }"""
            )

            total = await page.evaluate(
                "() => (typeof deckSlideCount === 'function' ? deckSlideCount() : document.querySelectorAll('.slide').length)"
            )
            if not total:
                raise SystemExit("未找到 .slide")

            for i in range(total):
                await page.evaluate(
                    """(i) => {
                      if (typeof deckGoTo === 'function') deckGoTo(i, true);
                      else if (typeof goTo === 'function') goTo(i);
                      else {
                        const slides = document.querySelectorAll('.slide');
                        slides.forEach((s, idx) => {
                          const on = idx === i;
                          s.classList.toggle('active', on);
                          s.style.display = on ? (s.classList.contains('slide-cover') ? 'grid' : 'flex') : 'none';
                        });
                      }
                    }""",
                    i,
                )
                await page.wait_for_timeout(args.settle_ms)
                fit = await page.evaluate(FIT_SLIDE_JS)
                # 缩放后再等一帧布局稳定
                await page.wait_for_timeout(120 if (fit or {}).get("scale", 1) < 0.999 else 60)

                shot = tmp_dir / f"slide_{i:03d}.jpg"
                # 截视口（背景色保留）；内容已 fit 进视口
                await page.screenshot(path=str(shot), type="jpeg", quality=92, full_page=False)
                shots.append(shot)

            await browser.close()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        for shot in shots:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(shot),
                Emu(0),
                Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

        prs.save(str(out_path))

    print(f"OK: {total} slides → {out_path}")
    return out_path


def main() -> None:
    args = _parse_args()
    asyncio.run(_export(args))


if __name__ == "__main__":
    main()
